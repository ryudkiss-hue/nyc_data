"""
Universal Exporter: Multi-format export system for visualizations and reports.

Supports 3 export formats for all 73 charts + 5 narratives:
1. PDF - ReportLab: Chart image + statistics table
2. CSV - Pandas: Data + statistics as CSV
3. Excel - openpyxl: Data sheet + statistics sheet with formatting

Usage:
    from app.services.universal_exporter import UniversalExporter
    exporter = UniversalExporter()

    # Export chart to PDF
    pdf_bytes = exporter.export_figure_to_pdf(
        figure=go.Figure(),
        title="Phase B: Moran's I Analysis",
        statistics={"Morans_I": 0.342, "p_value": 0.05}
    )

    # Export data to CSV
    csv_str = exporter.export_data_to_csv(
        df=pd.DataFrame(),
        title="Phase C: Distribution"
    )

    # Export to Excel
    xlsx_bytes = exporter.export_data_to_excel(
        df=pd.DataFrame(),
        narrative="Key insight text",
        title="Phase E: Decomposition"
    )
"""

import io
import logging
from datetime import datetime
from typing import Any, Optional

import pandas as pd
import plotly.graph_objects as go

logger = logging.getLogger(__name__)

class UniversalExporter:
    """Multi-format exporter for visualizations and data."""

    def __init__(self):
        """Initialize exporter with dependencies check."""
        self._check_dependencies()

    @staticmethod
    def _check_dependencies() -> None:
        """Verify required libraries are installed."""
        missing = []
        for lib in ["reportlab", "openpyxl"]:
            try:
                __import__(lib)
            except ImportError:
                missing.append(lib)

        if missing:
            logger.warning(f"Optional libraries missing: {missing}. Some exports may be limited.")

    def export_figure_to_pdf(
        self,
        figure: go.Figure,
        title: str,
        statistics: dict[str, Any],
        narrative: Optional[str] = None,
        table_df: Optional[pd.DataFrame] = None,
    ) -> Optional[bytes]:
        """
        Export Plotly figure to PDF with statistics table.

        Args:
            figure: Plotly Figure object
            title: Report title
            statistics: Dict of statistics to display
            narrative: Optional narrative text

        Returns:
            bytes: PDF file content (or None on error)

        Format:
            ┌──────────────────────────┐
            │ Title                    │
            │ [Chart Image]            │
            ├──────────────────────────┤
            │ Statistics               │
            │ Key1: Value1             │
            │ Key2: Value2             │
            │ ...                      │
            ├──────────────────────────┤
            │ Narrative (if provided)  │
            └──────────────────────────┘
        """
        try:
            from reportlab.lib import colors
            from reportlab.lib.pagesizes import A4, letter
            from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
            from reportlab.lib.units import inch
            from reportlab.pdfgen import canvas
            from reportlab.platypus import (
                PageBreak,
                Paragraph,
                SimpleDocTemplate,
                Spacer,
                Table,
                TableStyle,
            )
        except ImportError:
            logger.error("ReportLab not installed. Install with: pip install reportlab")
            return None

        try:
            # Convert figure to image (PNG). Non-fatal: a missing/failed renderer
            # (e.g. kaleido or a headless Chrome) must NOT prevent the PDF from
            # generating — we simply omit the chart and keep the stats/table.
            img_bytes = b""
            if figure is not None and getattr(figure, "data", None):
                try:
                    img_bytes = figure.to_image(format="png", width=700, height=500) or b""
                except Exception as e:  # noqa: BLE001 - renderer is optional
                    logger.warning(f"Chart image unavailable ({e}); exporting text-only PDF.")
                    img_bytes = b""

            # Create PDF buffer
            pdf_buffer = io.BytesIO()
            doc = SimpleDocTemplate(pdf_buffer, pagesize=letter)
            elements = []
            styles = getSampleStyleSheet()

            # Add title
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Heading1'],
                fontSize=24,
                textColor=colors.HexColor('#212529'),
                spaceAfter=30,
            )
            elements.append(Paragraph(title, title_style))

            # Add chart image
            if img_bytes:
                try:
                    from reportlab.platypus import Image
                    img_file = io.BytesIO(img_bytes)
                    img = Image(img_file, width=6*inch, height=4*inch)
                    elements.append(img)
                    elements.append(Spacer(1, 0.3*inch))
                except Exception as e:
                    logger.warning(f"Could not embed image: {e}")

            # Add statistics table
            if statistics:
                stat_data = [["Metric", "Value"]]
                for key, value in statistics.items():
                    if isinstance(value, float):
                        formatted = f"{value:.4f}"
                    else:
                        formatted = str(value)
                    stat_data.append([str(key), formatted])

                stat_table = Table(stat_data, colWidths=[3*inch, 3*inch])
                stat_table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#007bff')),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, 0), 12),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('BACKGROUND', (0, 1), (-1, -1), colors.HexColor('#f8f9fa')),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black),
                ]))
                elements.append(stat_table)

            # Add a data preview table (first rows/cols of the filtered data)
            if table_df is not None and not table_df.empty:
                elements.append(Spacer(1, 0.3 * inch))
                preview = table_df.head(15)
                cols = list(preview.columns)[:6]
                head = [str(c) for c in cols]
                body = [[str(v)[:24] for v in row]
                        for row in preview[cols].itertuples(index=False)]
                data_tbl = Table([head] + body)
                data_tbl.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#343a40')),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, -1), 7),
                    ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
                    ('ROWBACKGROUNDS', (0, 1), (-1, -1),
                     [colors.white, colors.HexColor('#f1f3f5')]),
                ]))
                elements.append(Paragraph(
                    f"<b>Data preview</b> ({len(table_df):,} rows × "
                    f"{len(table_df.columns)} cols)", styles['BodyText']))
                elements.append(Spacer(1, 0.1 * inch))
                elements.append(data_tbl)

            # Add narrative if provided
            if narrative:
                elements.append(Spacer(1, 0.2*inch))
                narrative_style = ParagraphStyle(
                    'Narrative',
                    parent=styles['BodyText'],
                    fontSize=10,
                    textColor=colors.HexColor('#495057'),
                )
                elements.append(Paragraph("<b>Key Insight</b>", narrative_style))
                elements.append(Paragraph(narrative[:500], narrative_style))

            # Build PDF
            doc.build(elements)
            pdf_buffer.seek(0)
            logger.info(f"PDF exported successfully: {title}")
            return pdf_buffer.getvalue()

        except Exception as e:
            logger.error(f"Error exporting to PDF: {e}", exc_info=True)
            return None

    def export_data_to_csv(
        self,
        df: pd.DataFrame,
        title: str,
        statistics: Optional[dict[str, Any]] = None,
    ) -> Optional[str]:
        """
        Export data to CSV format.

        Args:
            df: DataFrame to export
            title: Report title (included as comment)
            statistics: Optional statistics dict

        Returns:
            str: CSV content (or None on error)

        Format:
            # Report Title
            # Generated: 2026-06-11 14:30:00
            # Records: 100
            #
            column1,column2,column3,...
            value1,value2,value3,...
            ...
        """
        try:
            csv_buffer = io.StringIO()

            # Add header comments
            csv_buffer.write(f"# {title}\n")
            csv_buffer.write(f"# Generated: {datetime.now().isoformat()}\n")
            csv_buffer.write(f"# Records: {len(df)}\n")

            # Add statistics as comments
            if statistics:
                csv_buffer.write("#\n# Statistics:\n")
                for key, value in statistics.items():
                    csv_buffer.write(f"# {key}: {value}\n")

            csv_buffer.write("#\n")

            # Add data
            df.to_csv(csv_buffer, index=False)
            logger.info(f"CSV exported successfully: {title} ({len(df)} rows)")
            return csv_buffer.getvalue()

        except Exception as e:
            logger.error(f"Error exporting to CSV: {e}", exc_info=True)
            return None

    @staticmethod
    def _excel_safe(df: pd.DataFrame) -> pd.DataFrame:
        """Stringify cells openpyxl can't serialize (dict/list/tuple/set), e.g.
        Socrata 'location' point columns. Returns a copy; scalar columns untouched."""
        out = df.copy()
        for col in out.columns:
            s = out[col]
            if s.dtype == object:
                sample = s.dropna().head(20)
                if any(isinstance(v, (dict, list, tuple, set)) for v in sample):
                    s = s.apply(
                        lambda v: str(v) if isinstance(v, (dict, list, tuple, set)) else v)
            # openpyxl can't serialize pandas NA/NaT — coerce missing to None, but
            # only for columns that actually contain nulls (avoids an expensive
            # whole-frame astype on large exports).
            if s.isna().any():
                s = s.astype(object).where(s.notna(), None)
            out[col] = s
        return out

    def export_data_to_excel(
        self,
        df: pd.DataFrame,
        title: str,
        narrative: Optional[str] = None,
        statistics: Optional[dict[str, Any]] = None,
    ) -> Optional[bytes]:
        """
        Export data to Excel format with multiple sheets.

        Args:
            df: DataFrame to export
            title: Report title
            narrative: Optional narrative text (goes in separate sheet)
            statistics: Optional statistics dict (goes in separate sheet)

        Returns:
            bytes: XLSX file content (or None on error)

        Format:
            Sheet "Data": Full DataFrame
            Sheet "Statistics": Statistics table (if provided)
            Sheet "Narrative": Narrative text (if provided)
            Sheet "Summary": Title, date, record count
        """
        try:
            from openpyxl import Workbook
            from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
            from openpyxl.utils.dataframe import dataframe_to_rows
        except ImportError:
            logger.error("openpyxl not installed. Install with: pip install openpyxl")
            return None

        try:
            df = self._excel_safe(df)
            wb = Workbook()
            wb.remove(wb.active)  # Remove default sheet

            # Summary sheet
            ws_summary = wb.create_sheet("Summary", 0)
            ws_summary["A1"] = title
            ws_summary["A1"].font = Font(size=16, bold=True)
            ws_summary["A3"] = "Generated"
            ws_summary["B3"] = datetime.now().isoformat()
            ws_summary["A4"] = "Records"
            ws_summary["B4"] = len(df)

            # Data sheet
            ws_data = wb.create_sheet("Data", 1)
            for row_idx, row in enumerate(dataframe_to_rows(df, index=False, header=True), 1):
                for col_idx, value in enumerate(row, 1):
                    ws_data.cell(row=row_idx, column=col_idx, value=value)

            # Format header row
            header_fill = PatternFill(start_color="007bff", end_color="007bff", fill_type="solid")
            header_font = Font(color="ffffff", bold=True)
            for cell in ws_data[1]:
                cell.fill = header_fill
                cell.font = header_font

            # Statistics sheet (if provided)
            if statistics:
                ws_stats = wb.create_sheet("Statistics", 2)
                ws_stats["A1"] = "Metric"
                ws_stats["B1"] = "Value"
                for row_idx, (key, value) in enumerate(statistics.items(), 2):
                    ws_stats[f"A{row_idx}"] = key
                    ws_stats[f"B{row_idx}"] = value

            # Narrative sheet (if provided)
            if narrative:
                ws_narrative = wb.create_sheet("Narrative", 3)
                ws_narrative["A1"] = "Key Insight"
                ws_narrative["A2"] = narrative

            # Auto-size columns
            for ws in wb.sheetnames:
                ws_obj = wb[ws]
                for column in ws_obj.columns:
                    max_length = max(len(str(cell.value or "")) for cell in column)
                    ws_obj.column_dimensions[column[0].column_letter].width = min(max_length + 2, 50)

            # Export to bytes
            excel_buffer = io.BytesIO()
            wb.save(excel_buffer)
            excel_buffer.seek(0)
            logger.info(f"Excel exported successfully: {title}")
            return excel_buffer.getvalue()

        except Exception as e:
            logger.error(f"Error exporting to Excel: {e}", exc_info=True)
            return None

    def export_data_to_pptx(
        self,
        df: pd.DataFrame,
        title: str,
        statistics: Optional[dict[str, Any]] = None,
        narrative: Optional[str] = None,
        figure: Optional[go.Figure] = None,
    ) -> Optional[bytes]:
        """Export a slide deck: title + stats + data-preview table (+ chart if renderable).

        Uses python-pptx (pure-python). The chart image is best-effort — a missing
        renderer omits it rather than failing the whole export.
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError:
            logger.error("python-pptx not installed. Install with: pip install python-pptx")
            return None

        try:
            prs = Presentation()
            blank = prs.slide_layouts[6]

            # --- Title slide ---
            title_slide = prs.slides.add_slide(prs.slide_layouts[5])
            title_slide.shapes.title.text = title
            sub = title_slide.shapes.add_textbox(Inches(0.5), Inches(1.8),
                                                 Inches(9), Inches(1))
            tf = sub.text_frame
            tf.text = f"Generated {datetime.now().isoformat(timespec='seconds')}  •  {len(df):,} records"
            tf.paragraphs[0].font.size = Pt(14)
            if narrative:
                p = tf.add_paragraph()
                p.text = narrative[:300]
                p.font.size = Pt(12)

            # --- Statistics slide ---
            if statistics:
                s = prs.slides.add_slide(blank)
                box = s.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
                box.text_frame.text = "Key Statistics"
                box.text_frame.paragraphs[0].font.size = Pt(24)
                rows = len(statistics) + 1
                tbl = s.shapes.add_table(rows, 2, Inches(0.5), Inches(1.2),
                                         Inches(6), Inches(0.4 * rows)).table
                tbl.cell(0, 0).text, tbl.cell(0, 1).text = "Metric", "Value"
                for i, (k, v) in enumerate(statistics.items(), 1):
                    tbl.cell(i, 0).text = str(k)
                    tbl.cell(i, 1).text = f"{v:.4f}" if isinstance(v, float) else str(v)

            # --- Chart slide (best-effort) ---
            if figure is not None and getattr(figure, "data", None):
                try:
                    img = figure.to_image(format="png", width=900, height=560)
                    if img:
                        s = prs.slides.add_slide(blank)
                        s.shapes.add_picture(io.BytesIO(img), Inches(0.4), Inches(0.4),
                                             width=Inches(9))
                except Exception as e:  # noqa: BLE001 - renderer optional
                    logger.warning(f"PPTX chart image unavailable ({e}); skipping slide.")

            # --- Data preview slide ---
            if df is not None and not df.empty:
                s = prs.slides.add_slide(blank)
                box = s.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(9), Inches(0.5))
                box.text_frame.text = f"Data preview — {len(df):,} rows × {len(df.columns)} cols"
                box.text_frame.paragraphs[0].font.size = Pt(20)
                preview = df.head(10)
                cols = list(preview.columns)[:6]
                tbl = s.shapes.add_table(len(preview) + 1, len(cols), Inches(0.4),
                                         Inches(1.0), Inches(9), Inches(0.35 * (len(preview) + 1))).table
                for j, c in enumerate(cols):
                    tbl.cell(0, j).text = str(c)[:18]
                for i, row in enumerate(preview[cols].itertuples(index=False), 1):
                    for j, v in enumerate(row):
                        tbl.cell(i, j).text = str(v)[:20]

            buf = io.BytesIO()
            prs.save(buf)
            buf.seek(0)
            logger.info(f"PPTX exported successfully: {title}")
            return buf.getvalue()
        except Exception as e:
            logger.error(f"Error exporting to PPTX: {e}", exc_info=True)
            return None

    def export_report_multiformat(
        self,
        figure: go.Figure,
        df: pd.DataFrame,
        title: str,
        narrative: str,
        statistics: dict[str, Any],
    ) -> dict[str, Optional[bytes]]:
        """
        Export a complete report in all 3 formats.

        Args:
            figure: Plotly Figure
            df: DataFrame
            title: Report title
            narrative: Narrative text
            statistics: Statistics dict

        Returns:
            dict: {"pdf": bytes, "csv": str, "xlsx": bytes}

        Usage:
            exports = exporter.export_report_multiformat(
                figure=fig,
                df=df,
                title="Phase B Analysis",
                narrative="Key insight...",
                statistics={"morans_i": 0.342}
            )
            # Save to files
            with open("report.pdf", "wb") as f:
                f.write(exports["pdf"])
        """
        return {
            "pdf": self.export_figure_to_pdf(figure, title, statistics, narrative),
            "csv": self.export_data_to_csv(df, title, statistics),
            "xlsx": self.export_data_to_excel(df, title, narrative, statistics),
        }
