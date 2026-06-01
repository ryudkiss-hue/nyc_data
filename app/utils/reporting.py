"""Report generation utilities — PDF, PowerPoint, Excel."""

from __future__ import annotations

import io
from datetime import date
from pathlib import Path

import pandas as pd

try:
    import plotly.graph_objects as go  # noqa: F401

    HAS_PLOTLY = True
except ImportError:
    HAS_PLOTLY = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt  # noqa: F401

    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


# ---------------------------------------------------------------------------
# Item 77 — PDF export (WeasyPrint)
# ---------------------------------------------------------------------------


def dataframe_to_html_table(df: pd.DataFrame, max_rows: int = 100) -> str:
    """Render a DataFrame as an HTML table string."""
    return df.head(max_rows).to_html(index=False, border=1, classes="report-table")


def generate_pdf_report(
    title: str,
    summary_text: str,
    dataframes: dict[str, pd.DataFrame],
    output_bytes: bool = True,
) -> bytes | str:
    """Generate a PDF report from title, summary text, and named dataframes.

    Returns bytes if output_bytes=True, else saves to file and returns path.
    Requires WeasyPrint: pip install weasyprint
    """
    try:
        from weasyprint import HTML
    except ImportError:
        raise ImportError("WeasyPrint required: pip install weasyprint")

    tables_html = ""
    for name, df in dataframes.items():
        tables_html += f"<h2>{name}</h2>" + dataframe_to_html_table(df)

    html = f"""
    <!DOCTYPE html><html><head>
    <meta charset="utf-8">
    <style>
    body {{ font-family: Arial, sans-serif; margin: 2cm; }}
    h1 {{ color: #003087; border-bottom: 2px solid #003087; }}
    h2 {{ color: #003087; margin-top: 1.5em; }}
    .report-table {{ border-collapse: collapse; width: 100%; font-size: 10px; }}
    .report-table td, .report-table th {{ border: 1px solid #ddd; padding: 4px 6px; }}
    .report-table th {{ background: #003087; color: white; }}
    .summary {{ background: #f0f4f8; padding: 1em; border-radius: 4px; margin: 1em 0; }}
    </style>
    </head><body>
    <h1>🏙️ {title}</h1>
    <p class="summary">{summary_text}</p>
    {tables_html}
    <p style="font-size:9px;color:#999;margin-top:2em;">
        Generated: {date.today()} · NYC DOT SIM Toolkit
    </p>
    </body></html>
    """
    if output_bytes:
        return HTML(string=html).write_pdf()
    path = Path("data") / f"report_{date.today()}.pdf"
    path.parent.mkdir(exist_ok=True)
    HTML(string=html).write_pdf(str(path))
    return str(path)


# ---------------------------------------------------------------------------
# Item 78 — PowerPoint export
# ---------------------------------------------------------------------------


def generate_pptx_report(
    title: str,
    slides: list[dict],  # [{"title": str, "content": str, "df": pd.DataFrame | None}]
) -> bytes:
    """Generate a PPTX report. Each slide dict has title, content text, optional df table.

    Requires python-pptx: pip install python-pptx
    """
    if not HAS_PPTX:
        raise ImportError("python-pptx required: pip install python-pptx")

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Title slide
    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = f"Generated {date.today()} · NYC DOT SIM Toolkit"

    for s in slides:
        layout = prs.slide_layouts[1]
        sl = prs.slides.add_slide(layout)
        sl.shapes.title.text = s.get("title", "")
        if s.get("content"):
            sl.placeholders[1].text = s["content"]
        if s.get("df") is not None and not s["df"].empty:
            df = s["df"].head(20)
            rows, cols = len(df) + 1, len(df.columns)
            left, top, width, height = Inches(0.5), Inches(2.5), Inches(12), Inches(4)
            table = sl.shapes.add_table(rows, cols, left, top, width, height).table
            for ci, col in enumerate(df.columns):
                table.cell(0, ci).text = str(col)
            for ri, row in df.iterrows():
                for ci, val in enumerate(row):
                    table.cell(ri + 1, ci).text = str(val) if pd.notna(val) else ""

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# Item 79 — Scheduled email
# ---------------------------------------------------------------------------


def send_email_report(
    to_address: str,
    subject: str,
    body_text: str,
    attachment_bytes: bytes | None = None,
    attachment_name: str = "report.pdf",
    smtp_host: str = "localhost",
    smtp_port: int = 25,
) -> tuple[bool, str]:
    """Send an email with optional attachment via SMTP.

    Returns (success, message).
    """
    import smtplib
    from email.mime.application import MIMEApplication
    from email.mime.multipart import MIMEMultipart
    from email.mime.text import MIMEText

    msg = MIMEMultipart()
    msg["From"] = "nycdot-sim@localhost"
    msg["To"] = to_address
    msg["Subject"] = subject
    msg.attach(MIMEText(body_text, "plain"))

    if attachment_bytes:
        part = MIMEApplication(attachment_bytes, Name=attachment_name)
        part["Content-Disposition"] = f'attachment; filename="{attachment_name}"'
        msg.attach(part)

    try:
        with smtplib.SMTP(smtp_host, smtp_port, timeout=10) as server:
            server.sendmail(msg["From"], [to_address], msg.as_string())
        return True, f"Email sent to {to_address}"
    except Exception as e:
        return False, str(e)


# ---------------------------------------------------------------------------
# Item 80 — Excel multi-sheet export
# ---------------------------------------------------------------------------


def generate_excel_report(
    sheets: dict[str, pd.DataFrame],
    summary_text: str = "",
) -> bytes:
    """Generate multi-sheet Excel workbook. Keys become sheet names (truncated to 31 chars)."""
    buf = io.BytesIO()
    with pd.ExcelWriter(buf, engine="openpyxl") as writer:
        if summary_text:
            summary_df = pd.DataFrame([{"Summary": summary_text}])
            summary_df.to_excel(writer, sheet_name="Summary", index=False)
        for sheet_name, df in sheets.items():
            sheet = sheet_name[:31]
            df.to_excel(writer, sheet_name=sheet, index=False)
    return buf.getvalue()
