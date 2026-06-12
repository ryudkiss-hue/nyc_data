"""Publish destinations for Analyst Pack artifacts.

Goals:
- High-leverage automation that outputs concrete artifacts
- Dry-run support (returns an action plan without side effects)
- Minimal dependencies (PowerPoint export is optional via extras)
"""

from __future__ import annotations

import json
import os
import shutil
import smtplib
from dataclasses import dataclass
from email.message import EmailMessage
from pathlib import Path
from typing import Any

import yaml

from ..core.state import save_state


class PublishError(RuntimeError):
    pass

def _read_text_best_effort(path: Path) -> str:
    try:
        return path.read_text(encoding="utf-8")
    except Exception:
        try:
            return path.read_text(encoding="utf-8", errors="ignore")
        except Exception:
            return ""

def _pack_context(pack_dir: Path) -> dict[str, str]:
    manifest = pack_dir / "manifest.json"
    ctx: dict[str, str] = {
        "pack_name": pack_dir.name,
        "run_date": pack_dir.name,
        "profile_name": "",
    }
    if manifest.exists():
        try:
            data = json.loads(manifest.read_text(encoding="utf-8"))
            ctx["run_date"] = str(data.get("run_date") or ctx["run_date"])
            ctx["profile_name"] = str(data.get("profile_name") or "")
        except Exception:
            pass
    return ctx

def load_publish_profile(path: str | Path) -> dict[str, Any]:
    p = Path(path)
    if not p.exists():
        raise FileNotFoundError(f"Publish profile not found: {p}")
    data = yaml.safe_load(p.read_text(encoding="utf-8")) or {}
    if not isinstance(data, dict):
        raise PublishError("Publish profile must be a YAML mapping")
    data.setdefault("profile_name", p.stem)
    return data

@dataclass
class PublishAction:
    kind: str
    ok: bool
    detail: str
    meta: dict[str, Any]

@dataclass
class PublishReport:
    pack_dir: str
    profile_path: str
    dry_run: bool
    actions: list[PublishAction]

    def to_dict(self) -> dict[str, Any]:
        return {
            "pack_dir": self.pack_dir,
            "profile_path": self.profile_path,
            "dry_run": self.dry_run,
            "actions": [
                {"kind": a.kind, "ok": a.ok, "detail": a.detail, "meta": a.meta} for a in self.actions
            ],
        }

def _copy_pack(pack_dir: Path, dest_root: str, *, dry_run: bool) -> PublishAction:
    dest = Path(dest_root) / pack_dir.name
    if dry_run:
        return PublishAction(
            kind="file_copy",
            ok=True,
            detail=f"Would copy pack to {dest}",
            meta={"source": str(pack_dir), "dest": str(dest)},
        )
    try:
        dest.parent.mkdir(parents=True, exist_ok=True)
        if dest.exists():
            shutil.rmtree(dest)
        shutil.copytree(pack_dir, dest)
        return PublishAction(
            kind="file_copy",
            ok=True,
            detail=f"Copied pack to {dest}",
            meta={"source": str(pack_dir), "dest": str(dest)},
        )
    except Exception as exc:
        return PublishAction(
            kind="file_copy",
            ok=False,
            detail=f"Copy failed: {exc}",
            meta={"source": str(pack_dir), "dest": str(dest)},
        )

def _export_bi(pack_dir: Path, dest_root: str, include: list[str] | None, *, dry_run: bool) -> PublishAction:
    dest = Path(dest_root) / pack_dir.name
    wanted = set(include or [])
    candidates = sorted([p for p in pack_dir.iterdir() if p.is_file()])
    selected: list[Path] = []
    for p in candidates:
        if wanted and p.name not in wanted:
            continue
        if p.suffix.lower() in {".json", ".csv"}:
            selected.append(p)
        elif p.suffix.lower() in {".xlsx"}:
            selected.append(p)
        elif p.name.lower() in {"executive_summary.md", "executive_summary.html"}:
            selected.append(p)
    if dry_run:
        return PublishAction(
            kind="bi_export",
            ok=True,
            detail=f"Would export {len(selected)} files to {dest}",
            meta={"dest": str(dest), "files": [p.name for p in selected]},
        )
    dest.mkdir(parents=True, exist_ok=True)
    for p in selected:
        shutil.copy2(p, dest / p.name)
    return PublishAction(
        kind="bi_export",
        ok=True,
        detail=f"Exported {len(selected)} files to {dest}",
        meta={"dest": str(dest), "files": [p.name for p in selected]},
    )

def _teams_post(cfg: dict[str, Any], ctx: dict[str, str], summary: str, *, dry_run: bool) -> PublishAction:
    webhook = os.getenv(str(cfg.get("webhook_env", "TOOLKIT_TEAMS_WEBHOOK")), "")
    if not webhook:
        return PublishAction(
            kind="teams",
            ok=False,
            detail="Teams webhook not configured (set env var)",
            meta={"webhook_env": cfg.get("webhook_env", "TOOLKIT_TEAMS_WEBHOOK")},
        )
    title = str(cfg.get("title", "Analyst Pack published"))
    footer = str(cfg.get("footer", "")).strip()
    text = summary.strip()
    if footer:
        text = f"{text}\n\n{footer}"
    payload = {
        "@type": "MessageCard",
        "@context": "https://schema.org/extensions",
        "summary": title,
        "themeColor": "0078D7",
        "title": title,
        "sections": [{"facts": [{"name": "Pack", "value": ctx["pack_name"]}]}],
        "text": text[:9000],
    }
    if dry_run:
        return PublishAction(kind="teams", ok=True, detail="Would post to Teams webhook", meta={"payload": payload})
    import requests

    resp = requests.post(webhook, json=payload, timeout=20)
    ok = 200 <= resp.status_code < 300
    return PublishAction(
        kind="teams",
        ok=ok,
        detail=f"Teams webhook POST HTTP {resp.status_code}",
        meta={"status_code": resp.status_code, "response": resp.text[:500]},
    )

def _email_send(cfg: dict[str, Any], pack_dir: Path, ctx: dict[str, str], body_text: str, *, dry_run: bool) -> PublishAction:
    smtp_cfg = cfg.get("smtp") or {}
    host = str(smtp_cfg.get("host", "localhost"))
    port = int(smtp_cfg.get("port", 25))
    starttls = bool(smtp_cfg.get("starttls", False))
    username = os.getenv(str(smtp_cfg.get("username_env", "TOOLKIT_SMTP_USERNAME")), "")
    password = os.getenv(str(smtp_cfg.get("password_env", "TOOLKIT_SMTP_PASSWORD")), "")

    from_addr = os.getenv(str(cfg.get("from_env", "TOOLKIT_SMTP_FROM")), "") or username
    to_addrs = [str(x).strip() for x in (cfg.get("to") or []) if str(x).strip()]
    if not to_addrs:
        return PublishAction(kind="email", ok=False, detail="Email 'to' list empty", meta={})
    if not from_addr:
        return PublishAction(kind="email", ok=False, detail="Email from not configured", meta={})

    subject_tmpl = str(cfg.get("subject", "Analyst Pack — {pack_name}"))
    subject = subject_tmpl.format(**ctx)

    attach = [str(x) for x in (cfg.get("attach") or []) if str(x)]
    attach_paths: list[Path] = []
    for name in attach:
        p = pack_dir / name
        if p.exists() and p.is_file():
            attach_paths.append(p)

    msg = EmailMessage()
    msg["From"] = from_addr
    msg["To"] = ", ".join(to_addrs)
    msg["Subject"] = subject
    msg.set_content(body_text[:200000] if body_text else f"Analyst Pack: {ctx['pack_name']}")

    for p in attach_paths:
        data = p.read_bytes()
        maintype, subtype = ("application", "octet-stream")
        if p.suffix.lower() == ".html":
            maintype, subtype = ("text", "html")
        elif p.suffix.lower() == ".md":
            maintype, subtype = ("text", "markdown")
        elif p.suffix.lower() == ".json":
            maintype, subtype = ("application", "json")
        elif p.suffix.lower() == ".xlsx":
            maintype, subtype = ("application", "vnd.openxmlformats-officedocument.spreadsheetml.sheet")
        msg.add_attachment(data, maintype=maintype, subtype=subtype, filename=p.name)

    if dry_run:
        return PublishAction(
            kind="email",
            ok=True,
            detail=f"Would send email to {len(to_addrs)} recipient(s)",
            meta={"to": to_addrs, "subject": subject, "attachments": [p.name for p in attach_paths]},
        )

    try:
        with smtplib.SMTP(host=host, port=port, timeout=30) as smtp:
            smtp.ehlo()
            if starttls:
                smtp.starttls()
                smtp.ehlo()
            if username and password:
                smtp.login(username, password)
            smtp.send_message(msg)
        return PublishAction(
            kind="email",
            ok=True,
            detail=f"Sent email to {len(to_addrs)} recipient(s)",
            meta={"to": to_addrs, "subject": subject, "attachments": [p.name for p in attach_paths]},
        )
    except Exception as exc:
        return PublishAction(kind="email", ok=False, detail=f"Email send failed: {exc}", meta={})

def _pptx_export(cfg: dict[str, Any], pack_dir: Path, ctx: dict[str, str], executive_summary_text: str, *, dry_run: bool) -> PublishAction:
    try:
        from pptx import Presentation  # type: ignore
    except Exception as exc:
        return PublishAction(
            kind="pptx",
            ok=False,
            detail=f"python-pptx not installed (pip install '.[pptx]'): {exc}",
            meta={},
        )

    template_path = Path(str(cfg.get("template_path", "")))
    output_path = Path(str(cfg.get("output_path", f"outputs/published/{pack_dir.name}.pptx")))
    placeholders = dict(cfg.get("placeholders") or {})

    render_ctx = {
        **ctx,
        "executive_summary_text": executive_summary_text,
    }
    rendered = {k: str(v).format(**render_ctx) for k, v in placeholders.items()}

    if dry_run:
        return PublishAction(
            kind="pptx",
            ok=True,
            detail=f"Would render pptx to {output_path}",
            meta={"template": str(template_path), "output": str(output_path), "placeholders": rendered},
        )

    if not template_path.exists():
        return PublishAction(kind="pptx", ok=False, detail=f"Template not found: {template_path}", meta={})

    prs = Presentation(str(template_path))
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            tf = shape.text_frame
            for para in tf.paragraphs:
                for run in para.runs:
                    for token, replacement in rendered.items():
                        if token in run.text:
                            run.text = run.text.replace(token, replacement)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return PublishAction(kind="pptx", ok=True, detail=f"Wrote {output_path}", meta={"output": str(output_path)})

def publish_pack(
    *,
    pack_dir: str | Path,
    profile_path: str | Path,
    dry_run: bool = False,
    state_path: str = "outputs/.state/last_pack.json",
) -> PublishReport:
    pack = Path(pack_dir)
    if not pack.exists() or not pack.is_dir():
        raise FileNotFoundError(f"Pack directory not found: {pack}")

    profile = load_publish_profile(profile_path)
    ctx = _pack_context(pack)

    # Body source: best-effort based on preference list.
    body_prefer = [str(x) for x in ((profile.get("email") or {}).get("body_prefer") or [])]
    executive_summary_text = ""
    for name in body_prefer + ["executive_summary.md", "executive_summary.html"]:
        p = pack / name
        if p.exists():
            executive_summary_text = _read_text_best_effort(p)
            if executive_summary_text:
                break
    if not executive_summary_text:
        executive_summary_text = f"Analyst Pack: {ctx['pack_name']}"

    actions: list[PublishAction] = []

    def enabled(section: str) -> bool:
        cfg = profile.get(section) or {}
        return bool(cfg.get("enabled", False))

    if enabled("file_copy"):
        cfg = profile.get("file_copy") or {}
        dest_root = str(cfg.get("dest_root", "")).strip()
        if not dest_root:
            actions.append(PublishAction("file_copy", False, "file_copy.dest_root is required", meta={}))
        else:
            actions.append(_copy_pack(pack, dest_root, dry_run=dry_run))

    if enabled("bi_export"):
        cfg = profile.get("bi_export") or {}
        dest_root = str(cfg.get("dest_root", "")).strip()
        include = cfg.get("include")
        include_list = [str(x) for x in include] if isinstance(include, list) else None
        if not dest_root:
            actions.append(PublishAction("bi_export", False, "bi_export.dest_root is required", meta={}))
        else:
            actions.append(_export_bi(pack, dest_root, include_list, dry_run=dry_run))

    if enabled("teams"):
        actions.append(_teams_post(profile.get("teams") or {}, ctx, executive_summary_text, dry_run=dry_run))

    if enabled("email"):
        actions.append(_email_send(profile.get("email") or {}, pack, ctx, executive_summary_text, dry_run=dry_run))

    if enabled("pptx"):
        actions.append(_pptx_export(profile.get("pptx") or {}, pack, ctx, executive_summary_text, dry_run=dry_run))

    report = PublishReport(
        pack_dir=str(pack),
        profile_path=str(Path(profile_path)),
        dry_run=dry_run,
        actions=actions,
    )

    # Update state for Dash "Resume" UX (even on dry-run so UI knows user intent).
    try:
        # Multi-profile: if caller didn't override, mirror state into active profile folder too.
        from ..core.profiles import ensure_profile_exists

        save_state(
            state_path,
            {
                "last_pack_dir": str(pack),
                "last_publish_profile": str(Path(profile_path)),
                "last_publish_dry_run": bool(dry_run),
            },
        )
        try:
            prof = ensure_profile_exists()
            save_state(
                str(prof.state_dir / "last_pack.json"),
                {
                    "last_pack_dir": str(pack),
                    "last_publish_profile": str(Path(profile_path)),
                    "last_publish_dry_run": bool(dry_run),
                },
            )
        except Exception:
            pass
    except Exception:
        pass

    return report

