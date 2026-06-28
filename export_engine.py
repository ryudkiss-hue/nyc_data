import pandas as pd
import plotly.express as px
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader
import tempfile
from pathlib import Path
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO

DATASETS = {
    'Sidewalk Inspections': 'data/local_db/socrata_cache/inspection.parquet',
    'Violations Ledger': 'data/local_db/socrata_cache/violations.parquet',
    '311 Service Requests': 'data/local_db/socrata_cache/complaints_311.parquet',
    'Capital Projects': 'data/local_db/socrata_cache/capital_blocks.parquet'
}

def generate_pdf(datasets_to_export, row_limit, output_path):
    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    c = canvas.Canvas(output_path, pagesize=letter)
    width, height = letter

    # Title Page
    c.setFont("Helvetica-Bold", 24)
    c.drawString(50, height - 100, "Manhattan Mission Control")
    c.setFont("Helvetica", 16)
    c.drawString(50, height - 140, "Comprehensive Data Export (PDF)")
    c.setFont("Helvetica", 12)
    limit_str = f"{row_limit} rows" if row_limit > 0 else "No limit"
    c.drawString(50, height - 180, f"Data Limit: {limit_str}")
    c.showPage()

    for name in datasets_to_export:
        path = DATASETS.get(name)
        if not path or not os.path.exists(path):
            continue
            
        df = pd.read_parquet(path)
        if row_limit > 0:
            df = df.head(row_limit)
            
        y_pos = height - 50
        c.setFont("Helvetica-Bold", 18)
        c.drawString(50, y_pos, f"Dataset: {name}")
        y_pos -= 30
        
        c.setFont("Helvetica", 12)
        c.drawString(50, y_pos, f"Summary Statistics:")
        y_pos -= 20
        c.drawString(60, y_pos, f"Rows Exported: {len(df):,}")
        y_pos -= 20
        c.drawString(60, y_pos, f"Total Attributes: {len(df.columns)}")
        y_pos -= 30
        
        fig = _generate_fig(name, df)
        if fig:
            try:
                with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
                    fig.write_image(tmp.name, engine="kaleido", width=800, height=500)
                    img = ImageReader(tmp.name)
                    c.drawImage(img, 50, y_pos - 350, width=500, height=312)
                    y_pos -= 380
            except Exception as e:
                print(f"Failed to render image for {name}: {e}")
                
        c.showPage()
    c.save()

def generate_excel(datasets_to_export, row_limit, output_path):
    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    with pd.ExcelWriter(output_path, engine='xlsxwriter') as writer:
        for name in datasets_to_export:
            path = DATASETS.get(name)
            if not path or not os.path.exists(path):
                continue
            df = pd.read_parquet(path)
            if row_limit > 0:
                df = df.head(row_limit)
            sheet_name = name[:31]  # Excel sheet names max 31 chars
            df.to_excel(writer, sheet_name=sheet_name, index=False)

def generate_pptx(datasets_to_export, row_limit, output_path):
    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Manhattan Mission Control"
    limit_str = f"{row_limit} rows" if row_limit > 0 else "No limit"
    subtitle.text = f"Comprehensive Data Export\nData Limit: {limit_str}"

    for name in datasets_to_export:
        path = DATASETS.get(name)
        if not path or not os.path.exists(path):
            continue
            
        df = pd.read_parquet(path)
        if row_limit > 0:
            df = df.head(row_limit)
            
        # Slide with chart
        slide_layout = prs.slide_layouts[5] # Title only
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = f"Dashboard: {name} ({len(df):,} records)"
        
        fig = _generate_fig(name, df)
        if fig:
            try:
                img_stream = BytesIO()
                fig.write_image(img_stream, format='png', engine="kaleido", width=800, height=500)
                img_stream.seek(0)
                
                left = Inches(1)
                top = Inches(2)
                width = Inches(8)
                slide.shapes.add_picture(img_stream, left, top, width=width)
            except Exception as e:
                print(f"Failed to render image for {name} to PPTX: {e}")
                
    prs.save(output_path)

def _generate_fig(name, df):
    fig = None
    if name == 'Sidewalk Inspections' and 'noviolationfound' in df.columns:
        counts = df['noviolationfound'].value_counts().reset_index()
        fig = px.pie(counts, names='noviolationfound', values='count', title=f'{name} - Violation Found Distribution')
    elif name == 'Violations Ledger' and 'violationstatus' in df.columns:
        counts = df['violationstatus'].value_counts().reset_index()
        fig = px.bar(counts, x='violationstatus', y='count', title=f'{name} - Violation Status')
    elif name == '311 Service Requests' and 'status' in df.columns:
        counts = df['status'].value_counts().reset_index()
        fig = px.bar(counts, x='status', y='count', title=f'{name} - Complaint Status')
    elif name == 'Capital Projects' and 'project_type' in df.columns:
        counts = df['project_type'].value_counts().reset_index()
        fig = px.pie(counts, names='project_type', values='count', title=f'{name} - Project Types')
    return fig

def export_all(datasets, formats, row_limit):
    out_files = {}
    if "PDF" in formats:
        pdf_path = "docs/reports/export_report.pdf"
        generate_pdf(datasets, row_limit, pdf_path)
        out_files["PDF"] = pdf_path
    if "Excel" in formats:
        excel_path = "docs/reports/export_data.xlsx"
        generate_excel(datasets, row_limit, excel_path)
        out_files["Excel"] = excel_path
    if "PPTX" in formats:
        pptx_path = "docs/reports/export_presentation.pptx"
        generate_pptx(datasets, row_limit, pptx_path)
        out_files["PPTX"] = pptx_path
        
    return out_files
