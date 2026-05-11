"""BI & Presentation Integration for DOT Sidewalk Toolkit.

Export data and reports in formats compatible with:
- Tableau (TDS/Hyper via CSV + metadata, TDE compatible exports)
- Power BI (CSV with Power Query-friendly structure, DAX measure hints)
- PowerPoint (slide deck generation from reports and charts)
- Generic BI (structured JSON, pivot-ready CSVs)

Example::

    from socrata_toolkit.integrations.bi import (
        export_for_tableau,
        export_for_powerbi,
        create_presentation,
    )

    export_for_tableau(df, "output/tableau_data")
    export_for_powerbi(df, "output/powerbi_data")
    pptx_path = create_presentation(report, charts, "output/report.pptx")
"""

from __future__ import annotations

import csv
import json
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional, Sequence

import pandas as pd


# ---------------------------------------------------------------------------
# Tableau Export
# ---------------------------------------------------------------------------

@dataclass
class TableauExportResult:
    """Result from Tableau export operation."""
    csv_path: str
    metadata_path: str
    row_count: int
    column_count: int


def export_for_tableau(
    df: pd.DataFrame,
    output_dir: str,
    filename: str = "data",
    date_columns: Optional[List[str]] = None,
    geo_columns: Optional[Dict[str, str]] = None,
) -> TableauExportResult:
    """Export a DataFrame in a Tableau-optimized format.

    Creates:
    - A UTF-8 CSV file (Tableau's preferred format for text-based connections)
    - A metadata JSON file describing columns, types, and role hints

    Tableau can connect to this CSV directly. The metadata file provides
    field role hints (dimension vs measure, geographic roles) that can be
    used when configuring the data source.

    Args:
        df: Data to export.
        output_dir: Directory for output files.
        filename: Base filename (without extension).
        date_columns: Columns to format as ISO dates for Tableau parsing.
        geo_columns: Mapping of column name to Tableau geographic role
            (e.g., {"borough": "State/Province", "zipcode": "ZIP Code/Postcode"}).
    """
    out = Path(output_dir)
    out.mkdir(parents=True, exist_ok=True)
    csv_path = str(out / f"{filename}.csv")
    meta_path = str(out / f"{filename}_metadata.json")

    # Prepare data
    export_df = df.copy()
    if date_columns:
        for col in date_columns:
            if col in export_df.columns:
                export_df[col] = pd.to_datetime(export_df[col], errors="coerce").dt.strftime("%Y-%m-%d %H:%M:%S")

    export_df.to_csv(csv_path, index=False, quoting=csv.QUOTE_NONNUMERIC)

    # Build metadata
    columns_meta = []
    for col in df.columns:
        dtype = str(df[col].dtype)
        if "int" in dtype or "float" in dtype:
            role = "measure"
            data_type = "number"
        elif "datetime" in dtype or (date_columns and col in date_columns):
            role = "dimension"
            data_type = "datetime"
        else:
            role = "dimension"
            data_type = "string"

        geo_role = geo_columns.get(col) if geo_columns else None
        columns_meta.append({
            "name": col,
            "data_type": data_type,
            "role": role,
            "geographic_role": geo_role,
        })

    metadata = {
        "source": "socrata_toolkit",
        "row_count": len(df),
        "columns": columns_meta,
        "tableau_hints": {
            "connection_type": "textscan",
            "separator": ",",
            "text_qualifier": '"',
        },
    }
    Path(meta_path).write_text(json.dumps(metadata, indent=2), encoding="utf-8")

    return TableauExportResult(
        csv_path=csv_path,
        metadata_path=meta_path,
        row_count=len(df),
        column_count=len(df.columns),
    )


# ---------------------------------------------------------------------------
# Power BI Export
# ---------------------------------------------------------------------------

@dataclass
class PowerBIExportResult:
    """Result from Power BI export operation."""
    csv_path: str
    model_path: str
    dax_measures: List[Dict[str, str]]
    row_count: int


def export_for_powerbi(
    df: pd.DataFrame,
    output_dir: str,
    filename: str = "data",
    date_columns: Optional[List[str]] = None,
    measures: Optional[List[Dict[str, str]]] = None,
) -> PowerBIExportResult:
    """Export data in a Power BI-optimized format.

    Creates:
    - A CSV file suitable for Power Query import
    - A model JSON file with DAX measure definitions and column metadata

    Power BI can import the CSV, then the DAX measures can be copy-pasted
    into the model. The model file also includes suggested relationships
    and column formatting hints.

    Args:
        df: Data to export.
        output_dir: Directory for output files.
        filename: Base filename.
        date_columns: Columns to mark as date type.
        measures: List of DAX measure dicts with keys: name, expression, description.
    """
    out = Path(output_dir)
    out.mkdir(parents=True, exist_ok=True)
    csv_path = str(out / f"{filename}.csv")
    model_path = str(out / f"{filename}_model.json")

    export_df = df.copy()
    if date_columns:
        for col in date_columns:
            if col in export_df.columns:
                export_df[col] = pd.to_datetime(export_df[col], errors="coerce").dt.strftime("%Y-%m-%d")

    export_df.to_csv(csv_path, index=False)

    # Default DOT measures if none provided
    default_measures = measures or [
        {"name": "Total Violations", "expression": "SUM('data'[violations])", "description": "Sum of all violation counts"},
        {"name": "Avg Severity", "expression": "AVERAGE('data'[severity_rating])", "description": "Average severity rating"},
        {"name": "Pending Count", "expression": 'COUNTROWS(FILTER(\'data\', \'data\'[status] = "Pending Repair"))', "description": "Count of pending repairs"},
        {"name": "Completion Rate", "expression": 'DIVIDE(COUNTROWS(FILTER(\'data\', \'data\'[status] = "Complete")), COUNTROWS(\'data\'), 0)', "description": "Percentage of completed work"},
    ]

    columns_meta = []
    for col in df.columns:
        dtype = str(df[col].dtype)
        pbi_type = "Int64" if "int" in dtype else ("Double" if "float" in dtype else ("DateTime" if "datetime" in dtype else "String"))
        columns_meta.append({"name": col, "dataType": pbi_type})

    model = {
        "source": "socrata_toolkit",
        "table_name": filename,
        "columns": columns_meta,
        "measures": default_measures,
        "power_query_hint": f'let Source = Csv.Document(File.Contents("{filename}.csv"), [Delimiter=",", Encoding=65001, QuoteStyle=QuoteStyle.Csv]) in Source',
    }
    Path(model_path).write_text(json.dumps(model, indent=2), encoding="utf-8")

    return PowerBIExportResult(
        csv_path=csv_path,
        model_path=model_path,
        dax_measures=default_measures,
        row_count=len(df),
    )


# ---------------------------------------------------------------------------
# PowerPoint Presentation
# ---------------------------------------------------------------------------

@dataclass
class SlideContent:
    """Content for a single presentation slide."""
    title: str
    body: str
    data: Optional[Dict[str, Any]] = None
    chart_base64: Optional[str] = None
    layout: str = "content"  # "title", "content", "two_column", "chart"


def create_presentation(
    slides: List[SlideContent],
    output_path: str,
    title: str = "DOT Sidewalk Program Report",
    subtitle: str = "",
) -> str:
    """Create a PowerPoint-compatible presentation.

    If ``python-pptx`` is available, generates a .pptx file directly.
    Otherwise, falls back to generating an HTML presentation that can
    be opened in PowerPoint or Google Slides.

    Args:
        slides: List of SlideContent for each slide.
        output_path: Path for the output file.
        title: Presentation title (used on title slide).
        subtitle: Subtitle for title slide.

    Returns:
        The output file path.
    """
    p = Path(output_path)
    p.parent.mkdir(parents=True, exist_ok=True)

    try:
        return _create_pptx(slides, str(p), title, subtitle)
    except ImportError:
        return _create_html_presentation(slides, str(p), title, subtitle)


def _create_pptx(slides: List[SlideContent], path: str, title: str, subtitle: str) -> str:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if slide.placeholders[1]:
        slide.placeholders[1].text = subtitle

    # Content slides
    for sc in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = sc.title
        body = slide.placeholders[1]
        body.text = sc.body
        if sc.data:
            for key, val in sc.data.items():
                body.text += f"\n{key}: {val}"

    prs.save(path)
    return path


def _create_html_presentation(slides: List[SlideContent], path: str, title: str, subtitle: str) -> str:
    """Fallback: generate HTML slides importable by PowerPoint/Google Slides."""
    slide_html = []
    slide_html.append(f'<div class="slide title-slide"><h1>{title}</h1><h2>{subtitle}</h2></div>')

    for sc in slides:
        data_html = ""
        if sc.data:
            data_html = "<table>" + "".join(
                f"<tr><td><strong>{k}</strong></td><td>{v}</td></tr>" for k, v in sc.data.items()
            ) + "</table>"
        img_html = ""
        if sc.chart_base64:
            img_html = f'<img src="data:image/png;base64,{sc.chart_base64}" style="max-width:80%;margin:auto;display:block"/>'
        slide_html.append(f'<div class="slide"><h2>{sc.title}</h2><p>{sc.body}</p>{data_html}{img_html}</div>')

    html = f"""<!DOCTYPE html><html><head><meta charset="utf-8"><title>{title}</title>
<style>
body{{font-family:system-ui,sans-serif;margin:0;padding:0;background:#f0f0f0}}
.slide{{width:900px;min-height:500px;margin:40px auto;padding:60px;background:white;
box-shadow:0 2px 10px rgba(0,0,0,0.15);page-break-after:always}}
.title-slide{{text-align:center;padding-top:150px}}
h1{{color:#003366;font-size:36px}}h2{{color:#004488;font-size:24px}}
table{{border-collapse:collapse;margin:20px 0}}td{{padding:8px 16px;border:1px solid #ddd}}
</style></head><body>{"".join(slide_html)}</body></html>"""

    Path(path).write_text(html, encoding="utf-8")
    return path


# ---------------------------------------------------------------------------
# Generic BI Export (structured JSON for any platform)
# ---------------------------------------------------------------------------

def export_bi_package(
    df: pd.DataFrame,
    output_dir: str,
    dataset_name: str = "sidewalk_data",
    include_formats: Optional[List[str]] = None,
) -> Dict[str, str]:
    """Export a comprehensive BI package with multiple format options.

    Creates a directory with CSV, JSON, and metadata files suitable
    for import into any BI platform.

    Args:
        df: Source data.
        output_dir: Output directory.
        dataset_name: Name used for files.
        include_formats: List of formats to include. Default: all.

    Returns:
        Dict mapping format name to file path.
    """
    formats = include_formats or ["csv", "json", "parquet_schema"]
    out = Path(output_dir)
    out.mkdir(parents=True, exist_ok=True)
    paths: Dict[str, str] = {}

    if "csv" in formats:
        csv_path = str(out / f"{dataset_name}.csv")
        df.to_csv(csv_path, index=False)
        paths["csv"] = csv_path

    if "json" in formats:
        json_path = str(out / f"{dataset_name}.json")
        df.to_json(json_path, orient="records", indent=2, default_handler=str)
        paths["json"] = json_path

    # Schema metadata
    schema_path = str(out / f"{dataset_name}_schema.json")
    schema = {
        "dataset": dataset_name,
        "row_count": len(df),
        "columns": [
            {"name": col, "dtype": str(df[col].dtype), "null_count": int(df[col].isna().sum())}
            for col in df.columns
        ],
    }
    Path(schema_path).write_text(json.dumps(schema, indent=2), encoding="utf-8")
    paths["schema"] = schema_path

    return paths
