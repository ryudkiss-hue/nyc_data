"""Coverage tests for analyst.publish internal destination handlers."""

from __future__ import annotations

from unittest.mock import MagicMock, patch

import pytest


CTX = {"pack_name": "2024-01-01", "run_date": "2024-01-01"}


# ---------------------------------------------------------------------------
# _email_send
# ---------------------------------------------------------------------------

class TestEmailSend:
    def test_empty_to_list(self, tmp_path):
        from socrata_toolkit.analyst.publish import _email_send

        action = _email_send({"to": []}, tmp_path, CTX, "body", dry_run=False)
        assert action.ok is False
        assert "to' list empty" in action.detail

    def test_no_from(self, tmp_path, monkeypatch):
        from socrata_toolkit.analyst.publish import _email_send

        monkeypatch.delenv("TOOLKIT_SMTP_FROM", raising=False)
        monkeypatch.delenv("TOOLKIT_SMTP_USERNAME", raising=False)
        action = _email_send({"to": ["a@b.c"]}, tmp_path, CTX, "body", dry_run=False)
        assert action.ok is False
        assert "from not configured" in action.detail

    def test_dry_run(self, tmp_path, monkeypatch):
        from socrata_toolkit.analyst.publish import _email_send

        monkeypatch.setenv("TOOLKIT_SMTP_FROM", "from@x.com")
        action = _email_send({"to": ["a@b.c"]}, tmp_path, CTX, "body", dry_run=True)
        assert action.ok is True
        assert "Would send email" in action.detail

    def test_actual_send_with_attachment(self, tmp_path, monkeypatch):
        from socrata_toolkit.analyst.publish import _email_send

        monkeypatch.setenv("TOOLKIT_SMTP_FROM", "from@x.com")
        # create an attachment of each handled type
        (tmp_path / "report.html").write_text("<p>hi</p>", encoding="utf-8")
        (tmp_path / "data.json").write_text("{}", encoding="utf-8")
        cfg = {
            "to": ["a@b.c", "d@e.f"],
            "attach": ["report.html", "data.json", "missing.txt"],
            "smtp": {"host": "smtp.x.com", "port": 587, "starttls": True},
            "subject": "Pack {pack_name}",
        }
        smtp_instance = MagicMock()
        smtp_cm = MagicMock()
        smtp_cm.__enter__.return_value = smtp_instance
        smtp_cm.__exit__.return_value = False
        with patch("socrata_toolkit.analyst.publish.smtplib.SMTP", return_value=smtp_cm):
            action = _email_send(cfg, tmp_path, CTX, "body text", dry_run=False)
        assert action.ok is True
        assert "Sent email to 2" in action.detail
        smtp_instance.send_message.assert_called_once()

    def test_send_failure(self, tmp_path, monkeypatch):
        from socrata_toolkit.analyst.publish import _email_send

        monkeypatch.setenv("TOOLKIT_SMTP_FROM", "from@x.com")
        with patch("socrata_toolkit.analyst.publish.smtplib.SMTP", side_effect=OSError("conn refused")):
            action = _email_send({"to": ["a@b.c"]}, tmp_path, CTX, "body", dry_run=False)
        assert action.ok is False
        assert "Email send failed" in action.detail


# ---------------------------------------------------------------------------
# _pptx_export
# ---------------------------------------------------------------------------

class TestPptxExport:
    def test_pptx_not_installed(self, tmp_path):
        import sys

        from socrata_toolkit.analyst.publish import _pptx_export

        # Force the local `from pptx import Presentation` to fail.
        with patch.dict(sys.modules, {"pptx": None}):
            action = _pptx_export({}, tmp_path, CTX, "summary", dry_run=False)
        assert action.ok is False
        assert "python-pptx not installed" in action.detail

    def test_pptx_dry_run(self, tmp_path):
        import sys
        import types

        from socrata_toolkit.analyst.publish import _pptx_export

        fake_pptx = types.ModuleType("pptx")
        fake_pptx.Presentation = MagicMock()
        cfg = {
            "template_path": str(tmp_path / "template.pptx"),
            "output_path": str(tmp_path / "out.pptx"),
            "placeholders": {"{{title}}": "Report {pack_name}"},
        }
        with patch.dict(sys.modules, {"pptx": fake_pptx}):
            action = _pptx_export(cfg, tmp_path, CTX, "summary", dry_run=True)
        assert action.ok is True
        assert "Would render pptx" in action.detail

    def test_pptx_template_not_found(self, tmp_path):
        import sys
        import types

        from socrata_toolkit.analyst.publish import _pptx_export

        fake_pptx = types.ModuleType("pptx")
        fake_pptx.Presentation = MagicMock()
        cfg = {"template_path": str(tmp_path / "missing.pptx"), "output_path": str(tmp_path / "o.pptx")}
        with patch.dict(sys.modules, {"pptx": fake_pptx}):
            action = _pptx_export(cfg, tmp_path, CTX, "summary", dry_run=False)
        assert action.ok is False
        assert "Template not found" in action.detail


class TestPptxRender:
    def test_pptx_actual_render(self, tmp_path):
        pytest.importorskip("pptx")
        from pptx import Presentation

        from socrata_toolkit.analyst.publish import _pptx_export

        # Build a real template with a placeholder token
        template = tmp_path / "template.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        tf = slide.shapes.add_textbox(0, 0, 100, 50).text_frame
        tf.text = "{{title}}"
        prs.save(str(template))

        out = tmp_path / "rendered.pptx"
        cfg = {
            "template_path": str(template),
            "output_path": str(out),
            "placeholders": {"{{title}}": "Report {pack_name}"},
        }
        action = _pptx_export(cfg, tmp_path, {"pack_name": "2024-01-01"}, "summary", dry_run=False)
        assert action.ok is True
        assert out.exists()
        # Verify the placeholder was substituted
        rendered = Presentation(str(out))
        texts = [
            shape.text_frame.text
            for s in rendered.slides for shape in s.shapes
            if getattr(shape, "has_text_frame", False)
        ]
        assert any("Report 2024-01-01" in t for t in texts)
